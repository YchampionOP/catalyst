import os

import docx
import nltk
import pandas as pd
from pptx import Presentation
from PyPDF2 import PdfReader
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics import accuracy_score, classification_report
from sklearn.model_selection import train_test_split
from sklearn.naive_bayes import MultinomialNB

nltk.download('stopwords')
from nltk.corpus import stopwords


# Functions to extract text from different file types
def extract_text_from_docx(docx_file):
    doc = docx.Document(docx_file)
    return '\n'.join([para.text for para in doc.paragraphs])

def extract_text_from_pptx(pptx_file):
    prs = Presentation(pptx_file)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return '\n'.join(text_runs)

def extract_text_from_excel(excel_file):
    df = pd.read_excel(excel_file)
    return '\n'.join(df.astype(str).apply(lambda x: ' '.join(x), axis=1))

def extract_text_from_pdf(pdf_file):
    reader = PdfReader(pdf_file)
    text = ""
    for page in reader.pages:
        text += page.extract_text()
    return text

# Load and preprocess documents from various categories
def load_documents(directory):
    documents = []
    labels = []
    for category in os.listdir(directory):
        category_path = os.path.join(directory, category)
        if os.path.isdir(category_path):
            for file_name in os.listdir(category_path):
                file_path = os.path.join(category_path, file_name)
                file_extension = file_name.split('.')[-1]
                text = ""
                
                # Handle different file types
                if file_extension == 'txt':
                    with open(file_path, 'r', encoding='utf-8') as file:
                        text = file.read()
                elif file_extension == 'docx':
                    text = extract_text_from_docx(file_path)
                elif file_extension == 'pptx':
                    text = extract_text_from_pptx(file_path)
                elif file_extension == 'xlsx':
                    text = extract_text_from_excel(file_path)
                elif file_extension == 'pdf':
                    text = extract_text_from_pdf(file_path)
                
                if text:
                    documents.append(text)
                    labels.append(category)
                    
    return documents, labels

# Load documents from the 'Documents' folder
documents, labels = load_documents('Documents/')

# Preprocess text by lowercasing, removing stopwords, etc.
stop_words = set(stopwords.words('english'))

def preprocess_text(text):
    text = text.lower()
    text = ''.join([char for char in text if char.isalnum() or char.isspace()])
    words = text.split()
    words = [word for word in words if word not in stop_words]
    return ' '.join(words)

# Preprocess the loaded documents
processed_docs = [preprocess_text(doc) for doc in documents]

# Convert text to TF-IDF vectors
vectorizer = TfidfVectorizer()
X = vectorizer.fit_transform(processed_docs)

# Split the data into training and test sets (80% training, 20% testing)
X_train, X_test, y_train, y_test = train_test_split(X, labels, test_size=0.2, random_state=42)

# Train the Naive Bayes classifier
model = MultinomialNB()
model.fit(X_train, y_train)

# Test the model and print evaluation metrics
y_pred = model.predict(X_test)
print(f"Accuracy: {accuracy_score(y_test, y_pred)}")
print(classification_report(y_test, y_pred))

# Classify new documents
def classify_new_documents(new_docs):
    new_docs_processed = [preprocess_text(doc) for doc in new_docs]
    X_new = vectorizer.transform(new_docs_processed)
    predictions = model.predict(X_new)
    return predictions

# Example: Classify new documents
new_docs = ["The restaurant introduced a new menu with organic food options.", 
            "The movie was a blockbuster, making a huge profit in the first week."]
predicted_categories = classify_new_documents(new_docs)
for doc, category in zip(new_docs, predicted_categories):
    print(f"Document: '{doc}' -> Predicted Category: {category}")
