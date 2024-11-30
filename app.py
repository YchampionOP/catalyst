import os
import docx
import pandas as pd
from pptx import Presentation
from PyPDF2 import PdfReader
from flask import Flask, render_template, request, redirect, flash, url_for, session, send_from_directory
from werkzeug.utils import secure_filename
from flask_sqlalchemy import SQLAlchemy
from flask_bcrypt import Bcrypt
from flask_login import LoginManager, UserMixin, login_user, current_user, logout_user, login_required
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.naive_bayes import MultinomialNB
import nltk
import re
import requests
from flask_dance.contrib.github import make_github_blueprint, github
import base64
import json
from datetime import datetime

# Fetch user repositories from GitHub
def fetch_github_repos():
    headers = {'Authorization': f'token {app.config["GITHUB_TOKEN"]}'}
    response = requests.get(f'{app.config["GITHUB_API_URL"]}/user/repos', headers=headers)
    return response.json() if response.status_code == 200 else None

# Download a file from a GitHub repository
def download_github_file(repo, path):
    headers = {'Authorization': f'token {app.config["GITHUB_TOKEN"]}'}
    response = requests.get(f'{app.config["GITHUB_API_URL"]}/repos/{repo}/contents/{path}', headers=headers)
    if response.status_code == 200:
        file_content = response.json().get('content')
        return file_content  # Base64 encoded content
    return None


# Download stopwords
nltk.download('stopwords')
from nltk.corpus import stopwords

# Initialize Flask app
app = Flask(__name__)
app.config['UPLOAD_FOLDER'] = os.path.join('uploads', '{user_id}')  # Folder to save uploaded files
app.config['ALLOWED_EXTENSIONS'] = {'txt', 'pdf', 'docx', 'ppt', 'pptx', 'xls', 'xlsx'}
app.config['SECRET_KEY'] = 'e6f7b9dfb830d91de7b1da88b889954e' 
app.config['SQLALCHEMY_DATABASE_URI'] = 'sqlite:///site.db'
app.config['GITHUB_OAUTH_CLIENT_ID'] = 'Ov23liLWJvsCyio99dJ1'
app.config['GITHUB_OAUTH_CLIENT_SECRET'] = 'ee76e12893fd6cb794ce629f79d48d9fa8242b0d'
github_bp = make_github_blueprint(client_id=app.config['GITHUB_OAUTH_CLIENT_ID'],
                                    client_secret=app.config['GITHUB_OAUTH_CLIENT_SECRET'],
                                    redirect_to='github_integration',
                                    scope='repo')
app.register_blueprint(github_bp, url_prefix='/github_login')


# Initialize database, bcrypt, and login manager
db = SQLAlchemy(app)
bcrypt = Bcrypt(app)
login_manager = LoginManager(app)
login_manager.login_view = 'login'

# Ensure the upload folder exists
if not os.path.exists(app.config['UPLOAD_FOLDER']):
    os.makedirs(app.config['UPLOAD_FOLDER'])

class File(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    filename = db.Column(db.String(120), nullable=False)
    filepath = db.Column(db.String(255), nullable=False)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)

    user = db.relationship('User', backref=db.backref('files', lazy=True))

class Feedback(db.Model):
    id = db.Column(db.Integer, primary_key=True)
    user_id = db.Column(db.Integer, db.ForeignKey('user.id'), nullable=False)  # Link to the User table
    feedback = db.Column(db.Text, nullable=False)  # Feedback content
    timestamp = db.Column(db.DateTime, default=datetime.utcnow)  # Auto-generate timestamp

    user = db.relationship('User', backref=db.backref('feedbacks', lazy=True))  # Link feedback to user

# Initialize global variables for the vectorizer and models
vectorizer = TfidfVectorizer()
broad_category_model = MultinomialNB()  # Model for broad categories
subcategory_models = {}  # Dictionary to hold subcategory models

# Counter to track retraining iterations
retraining_iterations = 0
max_iterations = 10  # Limit the number of iterations

# User Model for authentication
class User(db.Model, UserMixin):
    id = db.Column(db.Integer, primary_key=True)
    username = db.Column(db.String(20), unique=True, nullable=False)
    email = db.Column(db.String(120), unique=True, nullable=False)
    password = db.Column(db.String(60), nullable=False)

@app.before_request
def enforce_https():
    if request.url.startswith('http://'):
        return redirect(request.url.replace('http://', 'https://', 1))

app.config['SESSION_COOKIE_SECURE'] = True
app.config['REMEMBER_COOKIE_SECURE'] = True

@app.after_request
def set_secure_headers(response):
    response.headers['Strict-Transport-Security'] = 'max-age=31536000; includeSubDomains'
    return response


@login_manager.user_loader
def load_user(user_id):
    return User.query.get(int(user_id))

def fetch_github_repos():
    if not github.authorized:
        return None
    resp = github.get('/user/repos')
    return resp.json() if resp.ok else None


# Helper function for username validation
def is_valid_username(username):
    return re.match("^[A-Za-z0-9_]+$", username) is not None

# Helper function for password validation
def is_valid_password(password):
    return (len(password) >= 8 and
            re.search(r"[A-Z]", password) and  # At least one uppercase letter
            re.search(r"[0-9]", password) and  # At least one number
            re.search(r"[!@#$%^&*(),.?\":{}|<>]", password))  # At least one special character

@app.route('/github.html', methods=['GET', 'POST'])
@login_required
def github_integration():
    if not github.authorized:
        flash('You are not logged into GitHub. Redirecting to GitHub login...', 'warning')
        return redirect(url_for('github.login'))
    
    # Fetch user repositories
    repos = fetch_github_repos()
    if repos is None:
        flash('Unable to fetch repositories from GitHub.', 'danger')
        return redirect(url_for('upload_and_classify'))

    # POST request handling
    if request.method == 'POST':
        repo_name = request.form['repo_name']
        file_path = request.form['file_path']
        file_content = download_github_file(repo_name, file_path)

        if file_content:
            # Decode base64 file content
            import base64
            decoded_content = base64.b64decode(file_content)
            new_docs = [decoded_content.decode('utf-8')]

            # Classify the document
            predictions = classify_new_documents(new_docs)
            return render_template('review.html', predictions=predictions)
        else:
            flash('Failed to download the file from GitHub.', 'danger')
            return redirect(url_for('github'))

    return render_template('github.html', repos=repos)

def fetch_github_repos():
    if not github.authorized:
        flash('GitHub authorization failed. Please log in again.', 'danger')
        return None

    # Make a request to GitHub's API to get the user's repositories
    resp = github.get('/user/repos')
    
    # Print the status and content for debugging purposes
    print(f"GitHub Response Status: {resp.status_code}")
    print(f"GitHub Response Content: {resp.text}")

    # Check if the response is OK and return repositories, else handle error
    if resp.ok:
        return resp.json()  # Return the repository JSON response
    else:
        flash('Failed to fetch repositories from GitHub.', 'danger')
        return None


# Routes for sign-up and login
@app.route('/signup', methods=['GET', 'POST'])
def signup():
    if current_user.is_authenticated:
        return redirect(url_for('upload_and_classify'))
    
    if request.method == 'POST':
        username = request.form['username']
        email = request.form['email']
        password = request.form['password']
        confirm_password = request.form['confirm_password']
        
        # Validate username
        if not is_valid_username(username):
            flash('Username can only contain letters, numbers, and underscores.', 'danger')
            return redirect(url_for('signup'))
        
        # Validate password
        if not is_valid_password(password):
            flash('Password must be at least 8 characters long, contain at least one uppercase letter, one number, and one special character.', 'danger')
            return redirect(url_for('signup'))
        
        # Check if the username or email already exists
        if User.query.filter((User.username == username) | (User.email == email)).first():
            flash('Username or Email already exists. Please choose a different one.', 'danger')
            return redirect(url_for('signup'))

        if password != confirm_password:
            flash('Passwords do not match.', 'danger')
            return redirect(url_for('signup'))

        hashed_password = bcrypt.generate_password_hash(password).decode('utf-8')
        user = User(username=username, email=email, password=hashed_password)
        db.session.add(user)
        db.session.commit()
        flash('Account created successfully! Please log in.', 'success')
        return redirect(url_for('login'))
    
    return render_template('signup.html')

@app.route('/login', methods=['GET', 'POST'])
def login():
    if current_user.is_authenticated:
        return redirect(url_for('upload_and_classify'))
   
    if request.method == 'POST':
        email = request.form['email']
        password = request.form['password']
        user = User.query.filter_by(email=email).first()
        if user and bcrypt.check_password_hash(user.password, password):
            login_user(user)
            # Create a user-specific directory if it doesn't exist
            user_folder = os.path.join('uploads', str(user.id))
            os.makedirs(user_folder, exist_ok=True)
            flash('Login successful!', 'success')
            return redirect(url_for('upload_and_classify'))
        else:
            flash('Login unsuccessful. Please check email and password.', 'danger')
    
    return render_template('login.html')

@app.route('/logout')
@login_required
def logout():
    logout_user()
    flash('You have been logged out.', 'info')
    return redirect(url_for('login'))

# Function to check if the file is allowed
def allowed_file(filename):
    return '.' in filename and filename.rsplit('.', 1)[1].lower() in app.config['ALLOWED_EXTENSIONS']

# Text preprocessing function
def preprocess_text(text):
    text = text.lower()  # Convert to lowercase
    text = ''.join([char for char in text if char.isalnum() or char.isspace()])  # Remove punctuation
    words = text.split()  # Tokenize
    words = [word for word in words if word not in stopwords.words('english')]  # Remove stopwords
    return ' '.join(words)

# Functions to extract text from different file types
def extract_text_from_docx(docx_file):
    doc = docx.Document(docx_file)
    return '\n'.join([para.text for para in doc.paragraphs])

def extract_text_from_pdf(pdf_file):
    reader = PdfReader(pdf_file)
    text = ""
    for page in reader.pages:
        text += page.extract_text()
    return text

# Load documents for both broad categories and subcategories
def load_documents(directory, is_subcategory=False):
    documents = []
    labels = []
    
    for category in os.listdir(directory):
        category_path = os.path.join(directory, category)
        if os.path.isdir(category_path):
            for file_name in os.listdir(category_path):
                file_path = os.path.join(category_path, file_name)
                file_extension = file_name.split('.')[-1]

                text = ""
                if file_extension == 'txt':
                    with open(file_path, 'r', encoding='utf-8') as f:
                        text = f.read()
                elif file_extension == 'pdf':
                    text = extract_text_from_pdf(file_path)
                elif file_extension == 'docx':
                    text = extract_text_from_docx(file_path)

                if text:
                    documents.append(text)
                    if is_subcategory:
                        labels.append(f"{category}")  # Label with the broad category
                    else:
                        labels.append(category)
    
    return documents, labels

# Train the model for both broad categories and subcategories
def train_model():
    global retraining_iterations
    if retraining_iterations >= max_iterations:
        print("Max iterations reached, skipping retraining.")
        return

    broad_docs, broad_labels = load_documents('Documents/')
    processed_broad_docs = [preprocess_text(doc) for doc in broad_docs]
    
    global vectorizer
    X_broad = vectorizer.fit_transform(processed_broad_docs)
    
    global broad_category_model
    broad_category_model.fit(X_broad, broad_labels)
    
    for broad_category in os.listdir('Sub-categories/'):
        sub_docs, sub_labels = load_documents(os.path.join('Sub-categories', broad_category), is_subcategory=True)
        processed_sub_docs = [preprocess_text(doc) for doc in sub_docs]
        
        X_sub = vectorizer.transform(processed_sub_docs)
        
        subcategory_model = MultinomialNB()
        subcategory_model.fit(X_sub, sub_labels)
        
        subcategory_models[broad_category] = subcategory_model

    retraining_iterations += 1
    print(f"Model retrained. Iteration: {retraining_iterations}")

# Classify new documents into broad categories and subcategories
def classify_new_documents(new_docs):
    new_docs_processed = [preprocess_text(doc) for doc in new_docs]
    X_new = vectorizer.transform(new_docs_processed)

    broad_category_predictions = broad_category_model.predict(X_new)

    final_predictions = []
    for i, broad_category in enumerate(broad_category_predictions):
        if broad_category in subcategory_models:
            subcategory_model = subcategory_models[broad_category]
            sub_X_new = X_new[i].reshape(1, -1)
            subcategory_prediction = subcategory_model.predict(sub_X_new)
            final_predictions.append(f"{broad_category}/{subcategory_prediction[0]}")
        else:
            final_predictions.append(broad_category)
    
    return final_predictions

def extract_text_from_ppt(ppt_file):
    ppt = Presentation(ppt_file)
    text = []
    for slide in ppt.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return '\n'.join(text)

def extract_text_from_excel(excel_file):
    df = pd.read_excel(excel_file, sheet_name=None)  # Load all sheets
    text = []
    for sheet_name, sheet_data in df.items():
        text.append(sheet_data.to_string())
    return '\n'.join(text)

# Extract text from .doc files (Windows environment with pywin32)
def extract_text_from_doc(doc_file):
    import win32com.client
    word = win32com.client.Dispatch("Word.Application")
    doc = word.Documents.Open(doc_file)
    text = doc.Content.Text
    doc.Close()
    word.Quit()
    return text

def get_file_structure(path):
    file_structure = {}
    for entry in os.scandir(path):
        if entry.is_dir():
            file_structure[entry.name] = {
                "type": "folder",
                "contents": get_file_structure(entry.path)
            }
        else:
            file_structure[entry.name] = {"type": "file"}
    return file_structure

# Route to handle file uploads and classification
@app.route('/', methods=['GET', 'POST'])
@login_required
def upload_and_classify():
    if request.method == 'POST':
        files = request.files.getlist('file')
        new_docs = []
        # User-specific upload folder
        user_folder = os.path.join(app.config['UPLOAD_FOLDER'].format(user_id=current_user.id))
        os.makedirs(user_folder, exist_ok=True)
        uploaded_file_paths = []  # To store paths of all uploaded files
        for file in files:
            if file and allowed_file(file.filename):
                filename = secure_filename(file.filename)
                
                # file path in UPLOAD_FOLDER
                file_path = os.path.join(app.config['UPLOAD_FOLDER'], filename)
                # User-specific file path
                user_file_path = os.path.join(user_folder, filename)
                # Save the file in both locations
                file.save(file_path)  # Save in the  folder
                file.save(user_file_path)    # Save in the user-specific folder
                uploaded_file_paths.append(user_file_path)
                # Extract text from the uploaded file
                file_extension = filename.split('.')[-1].lower()
                if file_extension == 'txt':
                    with open(file_path, 'r', encoding='utf-8') as f:
                        new_docs.append(f.read())
                elif file_extension == 'pdf':
                    new_docs.append(extract_text_from_pdf(file_path))
                elif file_extension == 'docx':
                    new_docs.append(extract_text_from_docx(file_path))
                elif file_extension == 'doc':
                    new_docs.append(extract_text_from_doc(file_path))
                elif file_extension in ['ppt', 'pptx']:
                    new_docs.append(extract_text_from_ppt(file_path))
                elif file_extension in ['xls', 'xlsx']:
                    new_docs.append(extract_text_from_excel(file_path))

        if new_docs:
            # Check if the vectorizer is fitted before classification
            if hasattr(vectorizer, "vocabulary_"):
                predictions = classify_new_documents(new_docs)
                zipped_data = zip(new_docs, predictions)
                # Organize files into folders based on predictions
                for file_path, prediction in zip(uploaded_file_paths, predictions):
                    broad_category, subcategory = prediction.split('/') if '/' in prediction else (prediction, None)
                    # Path to user's root folder
                    user_root =  os.path.join(app.config['UPLOAD_FOLDER'], str(current_user.id))
                    
                    # Create broad category folder
                    broad_folder = os.path.join(user_folder, broad_category)
                    os.makedirs(broad_folder, exist_ok=True)

                    # Create subcategory folder if it exists
                    final_folder = broad_folder
                    if subcategory:
                        sub_folder = os.path.join(broad_folder, subcategory)
                        final_folder = sub_folder

                    # Move the file to the appropriate folder
                    final_file_path = os.path.join(final_folder, os.path.basename(file_path))
                    os.rename(file_path, final_file_path)
                    flash('Files successfully classified and organized.', 'success')
                    
            else:
                flash('The model is not yet trained. Please wait.', 'danger')
                return redirect(url_for('upload_and_classify'))

            return render_template('review.html', zipped_data=zipped_data)
    
    return render_template('upload.html')

@app.route('/submit', methods=['POST'])
@login_required
def submit_review():
    docs = request.form.getlist('doc')
    labels = request.form.getlist('label')

    dataset_base_folder = 'Sub-categories/'
    broad_dataset_base_folder = 'Documents/'

    for doc, label in zip(docs, labels):
        broad_category, subcategory = label.split('/')
        category_folder = os.path.join(dataset_base_folder, broad_category, subcategory)
        
        if not os.path.exists(category_folder):
            os.makedirs(category_folder)

        file_path = os.path.join(category_folder, f'document_{hash(doc)}.txt')
        with open(file_path, 'w', encoding='utf-8') as f:
            f.write(doc)

        broad_category_folder = os.path.join(broad_dataset_base_folder, broad_category)
        if not os.path.exists(broad_category_folder):
            os.makedirs(broad_category_folder)

        broad_file_path = os.path.join(broad_category_folder, f'document_{hash(doc)}.txt')
        with open(broad_file_path, 'w', encoding='utf-8') as f:
            f.write(doc)

    flash('Documents successfully submitted!', 'success')
    return redirect(url_for('upload_and_classify'))

@app.route('/storage.html', methods=['GET'])
@login_required
def user_storage():
    user_folder = os.path.join('uploads', str(current_user.id))
    if not os.path.exists(user_folder):
        os.makedirs(user_folder)  # Ensure user folder exists
    file_structure = {"root": {"type": "folder", "contents": get_file_structure(user_folder)}}
    return render_template('storage.html', file_structure=file_structure)

@app.route('/create_folder', methods=['POST'])
@login_required
def create_folder():
    folder_name = request.form.get('folder_name')
    current_path = request.form.get('current_path')  # This will be a relative path

    if not folder_name:
        flash('Folder name cannot be empty.', 'danger')
        return redirect(url_for('user_storage'))

    user_folder = os.path.join('uploads', str(current_user.id), current_path)
    new_folder_path = os.path.join(user_folder, folder_name)

    try:
        os.makedirs(new_folder_path, exist_ok=True)
        flash(f'Folder "{folder_name}" created successfully!', 'success')
    except Exception as e:
        flash(f'Error creating folder: {e}', 'danger')

    return redirect(url_for('user_storage'))

@app.route('/upload_file', methods=['POST'])
@login_required
def upload_file():
    uploaded_file = request.files.get('file')
    current_path = request.form.get('current_path')  # Relative path to current folder

    if not uploaded_file:
        flash('No file selected.', 'danger')
        return redirect(url_for('user_storage'))

    user_folder = os.path.join('uploads', str(current_user.id), current_path)
    os.makedirs(user_folder, exist_ok=True)  # Ensure the directory exists

    file_path = os.path.join(user_folder, secure_filename(uploaded_file.filename))
    try:
        uploaded_file.save(file_path)
        flash(f'File "{uploaded_file.filename}" uploaded successfully!', 'success')
    except Exception as e:
        flash(f'Error uploading file: {e}', 'danger')

    return redirect(url_for('user_storage'))

@app.route('/help.html')
def show_help():
    return render_template('help.html')

@app.route('/update_profile', methods=['POST'])
@login_required
def update_profile():
    name = request.form.get('name')
    email = request.form.get('email')
    
    # Update user info in the database
    user = User.query.get(current_user.id)
    user.name = name
    user.email = email
    db.session.commit()
    
    flash("Profile updated successfully!", "success")
    return redirect(url_for('show_settings'))

@app.route('/change_password', methods=['POST'])
@login_required
def change_password():
    current_password = request.form.get('current_password')
    new_password = request.form.get('new_password')

    if bcrypt.check_password_hash(current_user.password, current_password):
        hashed_password = bcrypt.generate_password_hash(new_password).decode('utf-8')
        current_user.password = hashed_password
        db.session.commit()
        flash("Password updated successfully!", "success")
    else:
        flash("Incorrect current password.", "danger")

    return redirect(url_for('show_settings'))
@app.route('/settings.html')
def show_settings():
    return render_template('settings.html')

@app.route('/submit_feedback', methods=['POST'])
@login_required
def submit_feedback():
    feedback_text = request.form.get('feedback')

    # Save feedback in the database
    feedback = Feedback(user_id=current_user.id, feedback=feedback_text)
    db.session.add(feedback)
    db.session.commit()

    flash("Thank you for your feedback! It has been successfully submitted.", "success")
    return redirect(url_for('show_settings'))

@app.route('/download_file', methods=['GET'])
@login_required
def download_file():
    file_path = request.args.get('path')  # Relative file path
    user_folder = os.path.join('uploads', str(current_user.id))
    absolute_file_path = os.path.join(user_folder, file_path)

    if os.path.isfile(absolute_file_path):
        directory, filename = os.path.split(absolute_file_path)
        return send_from_directory(directory, filename, as_attachment=True)

    flash('File not found.', 'danger')
    return redirect(url_for('user_storage'))

with app.app_context():
    db.create_all()

if __name__ == '__main__':
    train_model()
    app.run(ssl_context=('cert.pem', 'key.pem'), debug=True)
